#!/usr/bin/env python3
"""
converter_core.py — HTML Presentation to PPTX Engine
Captures each slide as a 4K Retina screenshot via headless Chrome and assembles into PPTX.
Designed to be invoked by the native macOS HTML to PPTX Converter.app.
"""
import os
import sys
import subprocess
import tempfile
import shutil
import glob
import time

# Force unbuffered stdout for real-time log streaming to the Swift app
sys.stdout.reconfigure(line_buffering=True)
sys.stderr.reconfigure(line_buffering=True)

# Ensure environment PATH has node and standard paths
EXTRA_PATHS = [
    os.path.expanduser("~/.nvm/versions/node/v24.18.0/bin"),
    "/opt/homebrew/bin",
    "/usr/local/bin",
    "/usr/bin",
    "/bin",
    "/usr/sbin",
    "/sbin"
]
for p in glob.glob(os.path.expanduser("~/.nvm/versions/node/*/bin")):
    EXTRA_PATHS.insert(0, p)

current_path = os.environ.get("PATH", "")
os.environ["PATH"] = ":".join(EXTRA_PATHS) + ":" + current_path

def find_node_binary():
    node_which = shutil.which("node")
    if node_which and os.path.exists(node_which):
        return node_which
    candidates = glob.glob(os.path.expanduser("~/.nvm/versions/node/*/bin/node")) + [
        "/Users/arunthomas/.nvm/versions/node/v24.18.0/bin/node",
        "/opt/homebrew/bin/node",
        "/usr/local/bin/node",
        "/usr/bin/node"
    ]
    for c in candidates:
        if c and os.path.isfile(c) and os.access(c, os.X_OK):
            return c
    return "node"

NODE_BIN = find_node_binary()
CHROME_EXEC = "/Users/arunthomas/.cache/puppeteer/chrome-headless-shell/mac_arm-151.0.7922.47/chrome-headless-shell-mac-arm64/chrome-headless-shell"
NODE_MODULES = "/Users/arunthomas/.gemini/antigravity-ide/scratch/OmniSearch/node_modules"

CAPTURE_JS = """
const puppeteer = require('{node_modules}/puppeteer-core');
const path = require('path');
const fs = require('fs');

const htmlPath = process.argv[2];
const outputDir = process.argv[3];

async function capture() {
    const browser = await puppeteer.launch({
        executablePath: '{chrome_exec}',
        headless: 'new',
        args: ['--no-sandbox', '--disable-setuid-sandbox', '--disable-gpu']
    });

    const page = await browser.newPage();
    await page.setViewport({
        width: 1600,
        height: 900,
        deviceScaleFactor: 2.5
    });

    const fileUrl = 'file://' + path.resolve(htmlPath);
    await page.goto(fileUrl, { waitUntil: 'networkidle0', timeout: 30000 });

    // Let CSS animations and fonts settle
    await new Promise(r => setTimeout(r, 800));

    await page.evaluate(() => {
        const nav = document.getElementById('presentation-nav') || document.querySelector('nav');
        if (nav) nav.style.display = 'none';
        
        const deck = document.getElementById('presentation-deck') || document.querySelector('main');
        if (deck) {
            deck.style.paddingTop = '0px';
            deck.style.height = '100vh';
        }

        document.documentElement.style.fontSize = '20px';

        let slides = document.querySelectorAll('#presentation-deck > div');
        if (!slides || slides.length === 0) {
            slides = document.querySelectorAll('.slide, .slide-container');
        }

        slides.forEach(slide => {
            slide.style.height = '100vh';
            slide.style.minHeight = '100vh';
            slide.style.maxHeight = '100vh';
            slide.style.padding = '4.5rem 5% 2.5rem';
            slide.style.boxSizing = 'border-box';
        });

        document.querySelectorAll('.fade-up, .reveal, [class*="fade"]').forEach(el => {
            el.classList.add('visible');
            el.style.opacity = '1';
            el.style.transform = 'translateY(0)';
        });
    });

    const slideCount = await page.evaluate(() => {
        let s = document.querySelectorAll('#presentation-deck > div');
        if (!s || s.length === 0) s = document.querySelectorAll('.slide, .slide-container');
        return s.length;
    });

    if (!fs.existsSync(outputDir)) {
        fs.mkdirSync(outputDir, { recursive: true });
    } else {
        fs.readdirSync(outputDir).forEach(f => {
            if (f.endsWith('.png')) fs.unlinkSync(path.join(outputDir, f));
        });
    }

    console.log(`📸 Capturing ${slideCount} slides for: ${path.basename(htmlPath)}`);

    for (let i = 0; i < slideCount; i++) {
        await page.evaluate((idx) => {
            const deck = document.getElementById('presentation-deck') || document.querySelector('main');
            let slides = document.querySelectorAll('#presentation-deck > div');
            if (!slides || slides.length === 0) slides = document.querySelectorAll('.slide, .slide-container');
            if (deck && slides[idx]) {
                deck.scrollTop = slides[idx].offsetTop;
            }
        }, i);

        // Allow CSS transitions and lazy-loaded content to render
        await new Promise(r => setTimeout(r, 700));

        const imagePath = path.join(outputDir, `slide_${String(i + 1).padStart(3, '0')}.png`);
        await page.screenshot({ path: imagePath, type: 'png', fullPage: false });
        console.log(`  ✓ Slide ${i + 1}/${slideCount} captured in 4K`);
    }

    await browser.close();
}

capture().catch(err => {
    console.error("Capture Error:", err);
    process.exit(1);
});
""".replace("{chrome_exec}", CHROME_EXEC).replace("{node_modules}", NODE_MODULES)

def convert_single_html(html_file):
    html_file = os.path.abspath(html_file)
    if not os.path.exists(html_file):
        print(f"File not found: {html_file}")
        return False

    folder_dir = os.path.dirname(html_file)
    base_name = os.path.basename(html_file)
    name_without_ext = os.path.splitext(base_name)[0]
    out_name = f"{name_without_ext}.pptx"
    out_pptx = os.path.join(folder_dir, out_name)

    print(f"\n{'='*50}")
    print(f" Converting: {base_name}")
    print(f" Output:     {out_name}")
    print(f"{'='*50}")

    start_time = time.time()

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_js = os.path.join(temp_dir, "capture.js")
        temp_img_dir = os.path.join(temp_dir, "slides")
        os.makedirs(temp_img_dir, exist_ok=True)

        with open(temp_js, "w", encoding="utf-8") as f:
            f.write(CAPTURE_JS)

        cmd = [NODE_BIN, temp_js, html_file, temp_img_dir]
        res = subprocess.run(cmd, env=os.environ)
        if res.returncode != 0:
            print(f"❌ Failed capturing slides for {base_name}")
            return False

        # Assemble PPTX
        try:
            import pptx
            from pptx.util import Inches
        except ImportError:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
            import pptx
            from pptx.util import Inches

        prs = pptx.Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]

        images = sorted([f for f in os.listdir(temp_img_dir) if f.endswith('.png')])
        if not images:
            print(f"❌ No slides captured for {base_name}")
            return False

        for img in images:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(os.path.join(temp_img_dir, img), 0, 0, width=Inches(13.333), height=Inches(7.5))

        prs.save(out_pptx)
        elapsed = time.time() - start_time
        file_size_mb = os.path.getsize(out_pptx) / (1024 * 1024)
        print(f"\n🎉 SUCCESS! Created PowerPoint deck:")
        print(f"   📂 {out_pptx} ({len(images)} slides, {file_size_mb:.1f} MB, {elapsed:.1f}s)\n")
        return True

def choose_file_dialog():
    osa_cmd = """
    set chosen to choose file with prompt "Select an HTML presentation file to convert to PPTX:" of type {"html", "htm"}
    POSIX path of chosen
    """
    try:
        res = subprocess.run(["osascript", "-e", osa_cmd], capture_output=True, text=True)
        path = res.stdout.strip()
        return path if path else None
    except Exception:
        return None

def main():
    args = sys.argv[1:]
    targets = []

    if not args:
        print("\nDrag and drop an HTML file onto this script, or pick a file below:")
        selected = choose_file_dialog()
        if selected:
            targets.append(selected)
        else:
            print("No file selected. Exiting.")
            sys.exit(0)
    else:
        for arg in args:
            arg = os.path.abspath(arg.strip("'\""))
            if os.path.isdir(arg):
                for f in os.listdir(arg):
                    if f.endswith('.html') or f.endswith('.htm'):
                        targets.append(os.path.join(arg, f))
            elif os.path.isfile(arg):
                targets.append(arg)

    if not targets:
        print("No valid HTML files found to convert.")
        sys.exit(1)

    print(f"\nProcessing {len(targets)} HTML presentation(s)...")
    success_count = 0
    for t in targets:
        if convert_single_html(t):
            success_count += 1

    print(f"\n{'='*50}")
    print(f" Results: {success_count}/{len(targets)} presentations converted successfully")
    print(f"{'='*50}")

if __name__ == "__main__":
    main()
